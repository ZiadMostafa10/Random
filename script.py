import os
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Function to read Excel file and return worksheet names
def get_worksheet_names(file_path):
    try:
        # Read the Excel file
        excel_file = pd.ExcelFile(file_path)
        # Get the worksheet names
        worksheet_names = excel_file.sheet_names
        return worksheet_names
    except Exception as e:
        print(f"Error reading Excel file: {e}")
        return []

# Function to generate visualizations from raw data or pivot table data and save as PowerPoint presentation
def visualize_data(data, option, worksheet_name, prs=None):
    # Define visualization options
    options = {
        '1': {'type': 'line', 'title': 'Line Chart'},
        '2': {'type': 'bar', 'title': 'Bar Chart'},
        '3': {'type': 'pie', 'title': 'Pie Chart'},
        # Add more visualization options as needed
    }

    # If data is provided as a DataFrame, use it directly
    if isinstance(data, pd.DataFrame):
        df = data
    # If data is provided as a dictionary (pivot table), convert it to a DataFrame
    elif isinstance(data, dict):
        df = pd.DataFrame(data)
        df.set_index('Product Line', inplace=True)  # Set 'Product Line' as index

    plt.figure(figsize=(8, 6))

    if options[option]['type'] == 'line':
        for column in df.columns:
            plt.plot(df.index, df[column], marker='o', label=column)
        plt.xlabel('Product Line')  # Label for x-axis
        plt.ylabel('Values')  # Label for y-axis
        plt.grid()
        plt.title(options[option]['title'])
        plt.legend()
        plt.xticks(rotation=45)
    elif options[option]['type'] == 'bar':
        df.plot(kind='bar', figsize=(10, 6))
        plt.xlabel('Product Line')  # Label for x-axis
        plt.ylabel('Sum of Values')  # Label for y-axis
        plt.grid()
        plt.title(options[option]['title'])
        plt.xticks(rotation=0)  # Rotate x-axis labels for better readability
        plt.legend(title='Month')  # Add legend title for months
        plt.tight_layout()  # Adjust layout to prevent clipping of labels
    elif options[option]['type'] == 'pie':
        df.sum().plot(kind='pie', autopct='%1.1f%%', figsize=(8, 6))
        plt.title(options[option]['title'])

    # Save the visualization as an image
    img_filename = f'{worksheet_name}_{options[option]["type"]}.png'  # Filename with worksheet name and visualization type
    img_path = os.path.join('images', img_filename)  # Path to save the visualization image
    plt.savefig(img_path)

    # Show the visualization
    plt.show()

    # If prs is not provided, create a new presentation
    if prs is None:
        prs = Presentation()
    else:
        # If presentation exists, open it
        pass

    # Add a slide for the visualization
    slide_layout = prs.slide_layouts[5]  # Select layout for title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = options[option]['title']

    # Add the generated visualization to the slide
    left = Inches(1)
    top = Inches(2)
    slide.shapes.add_picture(img_path, left, top, width=Inches(7))

    # Save the presentation
    prs.save('presentation.pptx')

if __name__ == "__main__":
    # Prompt the user to input file path
    file_path = input("Enter the path to the Excel file: ")

    # Get the worksheet names
    worksheet_names = get_worksheet_names(file_path)

    # Display worksheet names
    if worksheet_names:
        print("Worksheet names:")
        for i, name in enumerate(worksheet_names, start=1):
            print(f"{i}. {name}")
        
        # Prompt the user to select a worksheet
        selected_index = int(input("Enter the index of the worksheet you want to work with: ")) - 1
        selected_worksheet = worksheet_names[selected_index]
        
        # Load data from the selected worksheet
        df = pd.read_excel(file_path, sheet_name=selected_worksheet)
        
        # Check if the PowerPoint presentation exists
        prs = Presentation('presentation.pptx') if os.path.exists('presentation.pptx') else None
        
        # Prompt the user to select a visualization option
        option = input("Enter the visualization option (1 for Line Chart, 2 for Bar Chart, 3 for Pie Chart): ")
        
        # Visualize data and save to PowerPoint
        visualize_data(df, option, selected_worksheet, prs)
    else:
        print("No worksheets found in the Excel file.")
